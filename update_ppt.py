from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Path to the PowerPoint file
ppt_path = 'Foodie-Finder-Overview-and-Purpose .pptx'
output_path = 'Foodie-Finder-Overview-and-Purpose-Updated.pptx'

# Check if the file exists
if not os.path.exists(ppt_path):
    print(f"File not found: {ppt_path}")
    exit(1)

# Load the presentation
try:
    prs = Presentation(ppt_path)
    print(f"Successfully loaded presentation with {len(prs.slides)} slides")
    
    # Update Slide 1 (Overview)
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if hasattr(shape, "text") and "Foodie Finder is a Flask-based web application" in shape.text:
            # Update the overview text
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = "Foodie Finder is a Flask-based web application that helps users discover recipes based on ingredients they already have. It uses machine learning (TF-IDF and cosine similarity) to match ingredients with suitable recipes, reducing food waste and inspiring culinary creativity."
            p.font.size = Pt(18)
    
    # Update Slide 2 (Features)
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if hasattr(shape, "text") and "Foodie Finder: Powerful Features" in shape.text:
            # The title is correct, now we need to add feature content
            continue
    
    # Clear existing content on slide 2 (except title)
    shapes_to_keep = []
    for shape in slide2.shapes:
        if hasattr(shape, "text") and "Foodie Finder: Powerful Features" in shape.text:
            shapes_to_keep.append(shape)
    
    # Store shapes to delete
    shapes_to_delete = [shape for shape in slide2.shapes if shape not in shapes_to_keep]
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Add new content box for features
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Add feature points
    features = [
        "Smart Ingredient Matching: Uses TF-IDF and cosine similarity to find recipes based on available ingredients",
        "Cooking Time Filters: Quick (0-30 min), Medium (30-60 min), or Long (60+ min) options",
        "User Authentication: Secure login/signup system with password hashing and session management",
        "Favorites System: Save and manage your favorite recipes",
        "User Profiles: Personalized user experience",
        "Responsive Design: Modern UI that works on all devices",
        "Recipe Details: View comprehensive recipe information including ingredients, instructions, and cooking time"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = "• " + feature
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Update Slide 4 (Future Steps)
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if hasattr(shape, "text") and "Future of Foodie Finder: Next Steps" in shape.text:
            # The title is correct
            continue
        
        if hasattr(shape, "text") and "Cloud Deployment:" in shape.text:
            # Update the future steps
            text_frame = shape.text_frame
            text_frame.clear()
            
            future_steps = [
                "Cloud Deployment: Migrate to AWS, Google Cloud, or Azure for better scalability",
                "Machine Learning Enhancements: Improve recommendation algorithm with user feedback",
                "Mobile App Development: Create native mobile apps for iOS and Android",
                "Social Features: Allow users to share recipes and cooking experiences",
                "Dietary Preferences: Add filtering by dietary restrictions (vegan, gluten-free, etc.)",
                "Ingredient Substitution: Suggest alternatives for missing ingredients",
                "Integration with Grocery Services: Connect with online grocery delivery"
            ]
            
            for step in future_steps:
                p = text_frame.add_paragraph()
                p.text = step
                p.font.size = Pt(18)
                p.space_after = Pt(12)
    
    # Save the updated presentation
    prs.save(output_path)
    print(f"Successfully updated presentation and saved as {output_path}")
    
except Exception as e:
    print(f"Error updating PowerPoint file: {e}")
