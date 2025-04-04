from pptx import Presentation
import os

# Path to the PowerPoint file
ppt_path = 'Foodie-Finder-Overview-and-Purpose .pptx'

# Check if the file exists
if not os.path.exists(ppt_path):
    print(f"File not found: {ppt_path}")
    exit(1)

# Load the presentation
try:
    prs = Presentation(ppt_path)
    print(f"Successfully loaded presentation with {len(prs.slides)} slides")
    
    # Print information about each slide
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        
        # Print the slide title if it exists
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text}")
        
        # Print text from all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Print only the first 100 characters of each text block to avoid overwhelming output
                text_preview = shape.text[:100] + "..." if len(shape.text) > 100 else shape.text
                print(f"Text: {text_preview}")
    
except Exception as e:
    print(f"Error reading PowerPoint file: {e}")
